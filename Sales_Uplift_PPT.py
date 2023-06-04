from pptx import Presentation
from pptx.util import Inches

# 创建PPT
presentation = Presentation()

# 添加标题页
slide = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "提升促销员销售保健品成交率的方法"
subtitle.text = "基于国际知名品牌和企业的成功经验"

# 添加概述页
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
slide.shapes.title.text = "概述"
content = slide.placeholders[1].text_frame
content.text = "本PPT将介绍如何提升促销员销售保健品的成交率。我们将参考国际知名品牌和企业的成功经验，并总结以下方法。"

# 添加培训与教育页
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
slide.shapes.title.text = "培训与教育"
content = slide.placeholders[1].text_frame
content.text = "2.1 提供全面的产品知识培训\n" \
               "2.2 建立健康保健品专业知识库\n" \
               "2.3 通过角色扮演和模拟销售情景培训销售技巧"

# 添加建立信任与关系页
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
slide.shapes.title.text = "建立信任与关系"
content = slide.placeholders[1].text_frame
content.text = "3.1 重视客户关系管理\n" \
               "3.2 建立长期合作伙伴关系\n" \
               "3.3 建立客户忠诚度计划"

# 添加个性化销售页
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
slide.shapes.title.text = "个性化销售"
content = slide.placeholders[1].text_frame
content.text = "4.1 根据客户需求定制销售方案\n" \
               "4.2 推荐适合客户的保健品组合\n" \
               "4.3 提供个性化的健康咨询和建议"

# 添加营销和推广页
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
slide.shapes.title.text = "营销和推广"
content = slide.placeholders[1].text_frame
content.text = "5.1 制定全面的营销计划\n" \
               "5.2 利用社交媒体和在线渠道进行推广\n" \
               "5.3 持续进行市场调研和竞争分析"

# 添加数据分析与反馈页
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
slide.shapes.title.text = "数据分析与反馈"
content = slide.placeholders[1].text_frame
content.text = "6.1 收集销售数据并进行分析\n" \
               "6.2 根据数据结果提供有针对性的反馈\n" \
               "6.3 持续改进和优化销售策略"

# 添加成功案例页
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
slide.shapes.title.text = "成功案例"
content = slide.placeholders[1].text_frame
content.text = "7.1 国际知名品牌A的销售策略\n" \
               "7.2 企业B的客户关系管理实践\n" \
               "7.3 品牌C在社交媒体上的推广案例"

# 添加总结页
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
slide.shapes.title.text = "总结"
content = slide.placeholders[1].text_frame
content.text = "8.1 培训与教育是提升促销员销售保健品成交率的基础\n" \
               "8.2 建立信任与关系能够增加客户忠诚度\n" \
               "8.3 个性化销售和营销推广是关键\n" \
               "8.4 数据分析和反馈可以持续改进销售策略\n" \
               "8.5 成功案例提供实践经验和借鉴"

# 保存PPT
presentation.save("sales_presentation.pptx")
